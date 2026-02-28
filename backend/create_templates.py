"""
Template Generator Script

Creates basic PowerPoint templates for the AI Presentation Generator.
Run this script to initialize the templates directory.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os


TEMPLATES_DIR = "templates"


def ensure_dir():
    os.makedirs(TEMPLATES_DIR, exist_ok=True)


def create_professional_template():
    """Create professional template."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Just save an empty presentation - styling is applied dynamically
    path = os.path.join(TEMPLATES_DIR, "professional.pptx")
    prs.save(path)
    print(f"Created: {path}")


def create_startup_template():
    """Create startup template."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    path = os.path.join(TEMPLATES_DIR, "startup.pptx")
    prs.save(path)
    print(f"Created: {path}")


def create_academic_template():
    """Create academic template."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    path = os.path.join(TEMPLATES_DIR, "academic.pptx")
    prs.save(path)
    print(f"Created: {path}")


def create_minimal_template():
    """Create minimal template."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    path = os.path.join(TEMPLATES_DIR, "minimal.pptx")
    prs.save(path)
    print(f"Created: {path}")


def main():
    print("Creating PowerPoint templates...")
    ensure_dir()
    
    create_professional_template()
    create_startup_template()
    create_academic_template()
    create_minimal_template()
    
    print("\nAll templates created successfully!")
    print(f"Templates location: {os.path.abspath(TEMPLATES_DIR)}")


if __name__ == "__main__":
    main()
