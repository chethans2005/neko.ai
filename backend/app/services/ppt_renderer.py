from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.services.theme_service import theme_service
from app.services.layout_service import layout_service, LayoutType


class PPTRenderer:
    """Renderer that produces styled PPTX presentations based on themes and layouts.

    Styling is deterministic, uses only colors, shapes and typography (no images).
    """

    FONT_NAME = 'Calibri'
    SLIDE_WIDTH_IN = 13.333
    SLIDE_HEIGHT_IN = 7.5

    def render(self, slides: list, theme_name: str = 'professional', title: str = None) -> Presentation:
        theme = theme_service.get_theme(theme_name)

        prs = Presentation()
        prs.slide_width = Inches(self.SLIDE_WIDTH_IN)
        prs.slide_height = Inches(self.SLIDE_HEIGHT_IN)

        # Optional title slide using session title
        if title:
            self._add_title_slide(prs, title, theme)

        # Render content slides
        for s in slides:
            layout = layout_service.choose_layout(s)
            if layout == LayoutType.TITLE_SLIDE:
                self._add_simple_title(prs, s, theme)
            elif layout == LayoutType.TWO_COLUMN:
                self._add_two_column(prs, s, theme)
            else:
                self._add_content_slide(prs, s, theme)

        return prs

    def _add_background(self, slide, theme: dict):
        bg_color = theme['background']
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(self.SLIDE_WIDTH_IN), Inches(self.SLIDE_HEIGHT_IN)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()
        # send to back
        spTree = slide.shapes._spTree
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def _add_accent_bar(self, slide, theme: dict, position: str = 'top'):
        # Deterministic accent: top horizontal bar
        accent = theme['accent']
        if position == 'top':
            h = 0.35
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(self.SLIDE_WIDTH_IN), Inches(h)
            )
        else:
            w = 0.45
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(w), Inches(self.SLIDE_HEIGHT_IN)
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = accent
        shape.line.fill.background()

    def _add_title_slide(self, prs: Presentation, title: str, theme: dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_background(slide, theme)
        self._add_accent_bar(slide, theme, position='top')

        # Centered title box
        box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(self.SLIDE_WIDTH_IN - 3.0), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.FONT_NAME
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = theme['title']
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub = slide.shapes.add_textbox(Inches(1.5), Inches(4.4), Inches(self.SLIDE_WIDTH_IN - 3.0), Inches(0.8))
        stf = sub.text_frame
        p2 = stf.paragraphs[0]
        p2.text = 'AI‑Generated Presentation'
        p2.font.name = self.FONT_NAME
        p2.font.size = Pt(20)
        p2.font.color.rgb = theme['text']
        p2.alignment = PP_ALIGN.CENTER

    def _add_simple_title(self, prs: Presentation, slide_data: dict, theme: dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_background(slide, theme)
        self._add_accent_bar(slide, theme, position='top')

        box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(self.SLIDE_WIDTH_IN - 2.0), Inches(2.0))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.name = self.FONT_NAME
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = theme['title']
        p.alignment = PP_ALIGN.CENTER

        # speaker notes
        notes = slide_data.get('speaker_notes')
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_content_slide(self, prs: Presentation, slide_data: dict, theme: dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_background(slide, theme)
        self._add_accent_bar(slide, theme, position='top')

        # Title at top
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(self.SLIDE_WIDTH_IN - 1.2), Inches(1.0))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.name = self.FONT_NAME
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = theme['title']

        # Content bullets
        content = slide_data.get('content') or []
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(self.SLIDE_WIDTH_IN - 1.6), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(0)
        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = self.FONT_NAME
            p.font.size = Pt(26)
            p.font.color.rgb = theme['text']
            p.space_before = Pt(6)
            p.space_after = Pt(6)

        # speaker notes
        notes = slide_data.get('speaker_notes')
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_two_column(self, prs: Presentation, slide_data: dict, theme: dict):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._add_background(slide, theme)
        self._add_accent_bar(slide, theme, position='top')

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(self.SLIDE_WIDTH_IN - 1.2), Inches(0.9))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get('title', '')
        p.font.name = self.FONT_NAME
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = theme['title']

        # Split content into two columns
        content = slide_data.get('content') or []
        half = (len(content) + 1) // 2
        left = content[:half]
        right = content[half:]

        left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches((self.SLIDE_WIDTH_IN - 1.8) / 2), Inches(5.2))
        left_tf = left_box.text_frame
        left_tf.word_wrap = True
        for i, item in enumerate(left):
            if i == 0:
                p = left_tf.paragraphs[0]
            else:
                p = left_tf.add_paragraph()
            p.text = item
            p.font.name = self.FONT_NAME
            p.font.size = Pt(26)
            p.font.color.rgb = theme['text']
            p.space_after = Pt(6)

        right_box = slide.shapes.add_textbox(Inches(0.6 + (self.SLIDE_WIDTH_IN - 1.8) / 2 + 0.2), Inches(1.6), Inches((self.SLIDE_WIDTH_IN - 1.8) / 2), Inches(5.2))
        right_tf = right_box.text_frame
        right_tf.word_wrap = True
        for i, item in enumerate(right):
            if i == 0:
                p = right_tf.paragraphs[0]
            else:
                p = right_tf.add_paragraph()
            p.text = item
            p.font.name = self.FONT_NAME
            p.font.size = Pt(26)
            p.font.color.rgb = theme['text']
            p.space_after = Pt(6)

        # speaker notes
        notes = slide_data.get('speaker_notes')
        if notes:
            slide.notes_slide.notes_text_frame.text = notes


ppt_renderer = PPTRenderer()
