"""
Template Service - PPT Template Management

Handles loading and managing PowerPoint templates.
"""
import os
from pathlib import Path
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.models.schemas import TemplateType


class TemplateService:
    """
    Service for managing PowerPoint templates.
    
    Handles:
    - Template loading and caching
    - Template metadata
    - Default styling per template
    """
    
    TEMPLATES_DIR = "templates"
    
    # Template configurations
    TEMPLATE_CONFIG = {
        TemplateType.PROFESSIONAL: {
            "name": "Professional",
            "description": "Clean, corporate design suitable for business presentations",
            "file": "professional.pptx",
            "title_color": RGBColor(0x1F, 0x49, 0x7D),  # Dark blue
            "accent_color": RGBColor(0x2E, 0x86, 0xAB),  # Light blue
            "background": RGBColor(0xFF, 0xFF, 0xFF),  # White,
        },
        TemplateType.STARTUP: {
            "name": "Startup",
            "description": "Modern, bold design for startup pitches and innovation",
            "file": "startup.pptx",
            "title_color": RGBColor(0xFF, 0x6B, 0x6B),  # Coral
            "accent_color": RGBColor(0x4E, 0xCB, 0xC4),  # Teal
            "background": RGBColor(0x2D, 0x3, 0x4E),  # Dark
        },
        TemplateType.ACADEMIC: {
            "name": "Academic",
            "description": "Scholarly design for research and educational presentations",
            "file": "academic.pptx",
            "title_color": RGBColor(0x6B, 0x2D, 0x5B),  # Purple
            "accent_color": RGBColor(0xE9, 0xA8, 0x20),  # Gold
            "background": RGBColor(0xF8, 0xF8, 0xF8),  # Light gray
        },
        TemplateType.MINIMAL: {
            "name": "Minimal",
            "description": "Simple, elegant design with focus on content",
            "file": "minimal.pptx",
            "title_color": RGBColor(0x33, 0x33, 0x33),  # Dark gray
            "accent_color": RGBColor(0x66, 0x66, 0x66),  # Medium gray
            "background": RGBColor(0xFF, 0xFF, 0xFF),  # White
        },
    }
    
    def __init__(self):
        self._ensure_templates_dir()
    
    def _ensure_templates_dir(self):
        """Ensure templates directory exists."""
        os.makedirs(self.TEMPLATES_DIR, exist_ok=True)
    
    def get_template_path(self, template_type: TemplateType) -> Optional[str]:
        """Get path to template file if it exists."""
        config = self.TEMPLATE_CONFIG.get(template_type)
        if not config:
            return None
        
        path = os.path.join(self.TEMPLATES_DIR, config["file"])
        if os.path.exists(path):
            return path
        
        return None
    
    def get_template_config(self, template_type: TemplateType) -> Dict:
        """Get template configuration."""
        return self.TEMPLATE_CONFIG.get(template_type, self.TEMPLATE_CONFIG[TemplateType.PROFESSIONAL])
    
    def list_templates(self) -> List[Dict]:
        """List all available templates with metadata."""
        templates = []
        for template_type, config in self.TEMPLATE_CONFIG.items():
            path = self.get_template_path(template_type)
            templates.append({
                "id": template_type.value,
                "name": config["name"],
                "description": config["description"],
                "available": path is not None
            })
        return templates
    
    def load_template(self, template_type: TemplateType) -> Presentation:
        """
        Load a PowerPoint template.
        
        Returns the template presentation or a new blank one if template not found.
        """
        path = self.get_template_path(template_type)
        
        if path:
            return Presentation(path)
        else:
            # Create new blank presentation
            return Presentation()
    
    def create_template(self, template_type: TemplateType):
        """
        Create a basic template file.
        
        This generates a simple template with proper styling.
        """
        prs = Presentation()
        
        # Set slide dimensions (16:9 widescreen)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        config = self.get_template_config(template_type)
        
        # Save template
        path = os.path.join(self.TEMPLATES_DIR, config["file"])
        prs.save(path)
        
        return path


# Global template service instance
template_service = TemplateService()
